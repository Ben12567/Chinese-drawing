from __future__ import annotations

from pathlib import Path

import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_ROOT = ROOT / "data/processed/clp4k_v6_union_en"
ASSET_ROOT = ROOT / "paper/ppt_assets/figure1"
OUTPUT_PATH = ROOT / "paper/editable_figure1_motivation_overview.pptx"


COLORS = {
    "blue": RGBColor(34, 82, 105),
    "red": RGBColor(157, 45, 39),
    "green": RGBColor(77, 105, 65),
    "green_dark": RGBColor(35, 89, 52),
    "cream": RGBColor(251, 248, 238),
    "paper": RGBColor(253, 251, 246),
    "border": RGBColor(210, 199, 184),
    "muted": RGBColor(105, 112, 120),
    "text": RGBColor(38, 44, 52),
    "light_red": RGBColor(254, 242, 240),
    "light_green": RGBColor(241, 248, 239),
}


def _asset(path: str) -> Path:
    return ROOT / path


PANEL_A_IMAGE = _asset("data/processed/clp4k_v6_union_en/images/clp4k_v6_union_en_000000.png")
PANEL_B_IMAGES = [
    _asset("outputs/paper_suite_v6/paper_lora_only_817/predictions_multiseed/clp4k_v6_union_en_000001__seed42.png"),
    _asset("outputs/paper_suite_strong/paper_controlnet_817_seed42/predictions_multiseed/clp4k_v6_union_en_000110__seed42.png"),
    _asset("outputs/paper_suite_strong/paper_ip_adapter_only_817_seed42/predictions_multiseed/clp4k_v6_union_en_000061__seed42.png"),
]
STYLE_REF = _asset("data/processed/clp4k_v6_union_en/images/clp4k_v6_union_en_000038.png")
OUTPUT_IMAGE = _asset("outputs/paper_suite_v6/paper_ours_817/predictions_multiseed/clp4k_v6_union_en_000001__seed42.png")
STRUCTURE_MAP = _asset("data/processed/clp4k_v6_union_en/structure_maps/clp4k_v6_union_en_000001.png")


def _prepare_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(10.8)
    return prs


def _cover_asset(source: Path, name: str, width: int, height: int) -> Path:
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    target = ASSET_ROOT / name
    image = Image.open(source).convert("RGB")
    iw, ih = image.size
    scale = max(width / iw, height / ih)
    resized = image.resize((int(iw * scale), int(ih * scale)), Image.Resampling.LANCZOS)
    left = max(0, (resized.width - width) // 2)
    top = max(0, (resized.height - height) // 2)
    cropped = resized.crop((left, top, left + width, top + height))
    cropped.save(target)
    return target


def _contain_asset(source: Path, name: str, width: int, height: int, fill=(255, 255, 255)) -> Path:
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    target = ASSET_ROOT / name
    image = Image.open(source).convert("RGB")
    image.thumbnail((width, height), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (width, height), fill)
    canvas.paste(image, ((width - image.width) // 2, (height - image.height) // 2))
    canvas.save(target)
    return target


def _add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 10,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    font: str = "Arial",
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["text"]
    return shape


def _add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    width: float = 0.8,
    rounded: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = COLORS["border"]
        shape.line.transparency = 100000
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    return shape


def _add_circle_label(slide, x: float, y: float, label: str, color: RGBColor):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.62), Inches(0.62))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = RGBColor(255, 255, 255)
    circle.line.width = Pt(1.0)
    _add_text(slide, x, y + 0.01, 0.62, 0.58, label, size=24, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def _add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, *, background: RGBColor | None = None):
    _add_rect(slide, x, y, w, h, fill=background or RGBColor(255, 255, 255), line=COLORS["border"], width=0.65, rounded=True)
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def _add_picture_exact(slide, path: Path, x: float, y: float, w: float, h: float, *, background: RGBColor | None = None):
    _add_rect(slide, x, y, w, h, fill=background or RGBColor(255, 255, 255), line=COLORS["border"], width=0.65, rounded=True)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def _add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: float = 1.5, arrow: bool = False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def _add_dashed_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor):
    line = _add_line(slide, x1, y1, x2, y2, color, width=1.2, arrow=True)
    line.line.dash_style = 4
    return line


def _add_warning(slide, x: float, y: float, number: str, title: str, desc: str, image_paths: list[Path]):
    _add_rect(slide, x, y, 5.00, 2.15, fill=COLORS["paper"], line=RGBColor(222, 179, 174), width=0.8, rounded=True)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 0.18), Inches(y + 0.20), Inches(0.32), Inches(0.32))
    tri.fill.solid()
    tri.fill.fore_color.rgb = COLORS["red"]
    tri.line.color.rgb = COLORS["red"]
    _add_text(slide, x + 0.56, y + 0.12, 3.9, 0.26, f"{number}. {title}", size=10.0, bold=True, color=COLORS["red"])
    _add_text(slide, x + 0.56, y + 0.42, 3.9, 0.30, desc, size=7.8, color=COLORS["text"])
    _add_text(slide, x + 4.60, y + 0.17, 0.25, 0.24, "×", size=19, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)
    img_y = y + 0.80
    if len(image_paths) == 1:
        asset = _cover_asset(image_paths[0], f"warning_{number}_wide.png", 900, 215)
        _add_picture_exact(slide, asset, x + 0.28, img_y, 4.44, 1.05, background=RGBColor(255, 255, 255))
    else:
        each = 1.36
        for idx, path in enumerate(image_paths):
            asset = _cover_asset(path, f"warning_{number}_{idx}.png", 280, 215)
            _add_picture_exact(slide, asset, x + 0.28 + idx * (each + 0.12), img_y, each, 1.05, background=RGBColor(255, 255, 255))


def _prepare_structure_assets() -> dict[str, Path]:
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    rgba = Image.open(STRUCTURE_MAP).convert("RGBA")
    arr = np.asarray(rgba)
    names = ["lineart", "quantized_depth", "blank_space_mask", "saliency_mask"]
    result: dict[str, Path] = {}
    for idx, name in enumerate(names):
        channel = Image.fromarray(arr[..., idx]).convert("L")
        if name == "lineart":
            channel = ImageOps.autocontrast(channel)
        else:
            channel = ImageOps.autocontrast(channel)
        raw = ASSET_ROOT / f"figure1_{name}_raw.png"
        channel.save(raw)
        out = _contain_asset(raw, f"figure1_{name}.png", 260, 76)
        raw.unlink(missing_ok=True)
        result[name] = out
    return result


def _add_panel_header(slide, x: float, y: float, w: float, label: str, title: str, color: RGBColor):
    _add_rect(slide, x, y, w, 0.70, fill=color, line=color, width=0.0, rounded=True)
    _add_circle_label(slide, x + 0.12, y + 0.10, label, color)
    _add_text(slide, x + 0.84, y + 0.09, w - 1.0, 0.50, title, size=14.5, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def _add_panel_a(slide):
    x, y, w, h = 0.18, 0.18, 6.95, 10.35
    _add_rect(slide, x, y, w, h, fill=COLORS["paper"], line=RGBColor(163, 193, 204), width=1.0, rounded=True)
    _add_panel_header(slide, x, y, w, "A", "Chinese landscape painting as a\nstructured cultural visual language", COLORS["blue"])
    panel_a_asset = _cover_asset(PANEL_A_IMAGE, "figure1_panel_a_painting.png", 820, 1670)
    _add_picture_exact(slide, panel_a_asset, x + 0.25, y + 1.05, 4.10, 8.35, background=RGBColor(255, 255, 255))

    _add_dashed_line(slide, x + 0.30, y + 3.25, x + 4.30, y + 3.25, RGBColor(89, 132, 151))
    _add_dashed_line(slide, x + 0.30, y + 5.62, x + 4.30, y + 5.62, RGBColor(89, 132, 151))
    _add_text(slide, x + 0.42, y + 1.70, 1.10, 0.36, "Background", size=8.6, bold=True, color=COLORS["blue"])
    _add_text(slide, x + 0.42, y + 4.20, 1.10, 0.36, "Midground", size=8.6, bold=True, color=COLORS["green_dark"])
    _add_text(slide, x + 0.42, y + 7.46, 1.10, 0.36, "Foreground", size=8.6, bold=True, color=RGBColor(150, 105, 24))

    features = [
        ("spatial\nhierarchy", "clear depth ordering\nof mountains, water,\nand objects.", COLORS["blue"]),
        ("blank\nspace", "intentional emptiness\ncreates breath and\nmeaning.", COLORS["green_dark"]),
        ("brush-and-ink\nrhythm", "varied brushwork,\nink intensity, and\ntextures convey vitality.", COLORS["text"]),
        ("poetic\natmosphere", "evokes emotion and\nliterary associations.", RGBColor(174, 100, 42)),
        ("composition", "arrangement follows\nprinciples of balance\nand harmony.", RGBColor(122, 91, 40)),
    ]
    fy = y + 1.22
    for idx, (title, desc, color) in enumerate(features):
        yy = fy + idx * 1.55
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 4.65), Inches(yy), Inches(0.62), Inches(0.62))
        icon.fill.background()
        icon.line.color.rgb = color
        icon.line.width = Pt(1.1)
        if idx == 0:
            _add_text(slide, x + 4.79, yy + 0.14, 0.33, 0.25, "山", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        elif idx == 1:
            _add_text(slide, x + 4.78, yy + 0.13, 0.34, 0.26, "○", size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        elif idx == 2:
            _add_text(slide, x + 4.78, yy + 0.11, 0.34, 0.30, "墨", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        elif idx == 3:
            _add_text(slide, x + 4.78, yy + 0.12, 0.34, 0.30, "景", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, x + 4.78, yy + 0.12, 0.34, 0.30, "法", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_text(slide, x + 5.40, yy - 0.02, 1.35, 0.45, title, size=10.0, bold=True, color=color)
        _add_text(slide, x + 5.40, yy + 0.44, 1.35, 0.76, desc, size=6.9, color=COLORS["text"], valign=MSO_ANCHOR.TOP)
        _add_dashed_line(slide, x + 4.18, y + 2.45 + idx * 1.48, x + 4.65, yy + 0.31, color)


def _add_panel_b(slide):
    x, y, w, h = 7.28, 0.18, 5.35, 10.35
    _add_rect(slide, x, y, w, h, fill=COLORS["light_red"], line=RGBColor(223, 169, 165), width=1.0, rounded=True)
    _add_panel_header(slide, x, y, w, "B", "Why generic diffusion models\nare insufficient", COLORS["red"])
    _add_warning(slide, x + 0.22, y + 0.98, "1", "structure confusion", "foreground and background are mixed.", [PANEL_B_IMAGES[0]])
    _add_warning(slide, x + 0.22, y + 3.30, "2", "blank space treated as background", "empty regions are either overfilled or awkward.", [PANEL_B_IMAGES[1], PANEL_B_IMAGES[2]])
    _add_warning(slide, x + 0.22, y + 5.62, "3", "style drift", "inconsistent brushwork and mixed styles.", [PANEL_B_IMAGES[0], PANEL_B_IMAGES[1], PANEL_B_IMAGES[2]])
    _add_rect(slide, x + 0.32, y + 8.48, w - 0.64, 1.22, fill=RGBColor(255, 247, 244), line=RGBColor(222, 179, 174), width=0.8, rounded=True)
    _add_circle_label(slide, x + 0.52, y + 8.72, "!", COLORS["red"])
    _add_text(
        slide,
        x + 1.18,
        y + 8.58,
        w - 1.55,
        0.94,
        "Generic text-to-image models may mimic surface style,\n"
        "but often fail to preserve composition,\n"
        "blank-space organization, and stylistic consistency.",
        size=8.7,
        bold=True,
        color=COLORS["red"],
        valign=MSO_ANCHOR.TOP,
    )


def _add_structure_table(slide, x: float, y: float, assets: dict[str, Path]):
    labels = ["lineart", "quantized\ndepth", "blank-space\nmask", "saliency\nmask"]
    keys = ["lineart", "quantized_depth", "blank_space_mask", "saliency_mask"]
    row_h = 0.48
    _add_rect(slide, x, y, 2.45, row_h * 4, fill=RGBColor(255, 255, 255), line=COLORS["border"], width=0.6, rounded=True)
    for idx, (label, key) in enumerate(zip(labels, keys)):
        yy = y + idx * row_h
        _add_rect(slide, x, yy, 0.82, row_h, fill=RGBColor(255, 255, 255), line=COLORS["border"], width=0.45)
        _add_text(slide, x + 0.05, yy + 0.05, 0.72, row_h - 0.10, label, size=5.7, color=COLORS["text"], align=PP_ALIGN.CENTER)
        slide.shapes.add_picture(str(assets[key]), Inches(x + 0.82), Inches(yy), Inches(1.63), Inches(row_h))


def _add_panel_c(slide):
    x, y, w, h = 12.82, 0.18, 7.00, 10.35
    _add_rect(slide, x, y, w, h, fill=COLORS["light_green"], line=RGBColor(183, 202, 171), width=1.0, rounded=True)
    _add_panel_header(slide, x, y, w, "C", "Our culture-aware hierarchical\ndiffusion framework", COLORS["green"])
    assets = _prepare_structure_assets()

    in_x = x + 0.26
    box_w = 2.75
    _add_rect(slide, in_x, y + 1.05, box_w, 1.52, fill=COLORS["paper"], line=COLORS["border"], width=0.65, rounded=True)
    _add_text(slide, in_x + 0.12, y + 1.12, box_w - 0.24, 0.22, "1. Structured text prompt", size=8.8, bold=True, color=COLORS["green_dark"])
    _add_text(
        slide,
        in_x + 0.16,
        y + 1.43,
        box_w - 0.32,
        0.92,
        "Misty mountains, river,\npavilion, poetic blank\nspace, ink-wash style",
        size=6.7,
        color=COLORS["text"],
        valign=MSO_ANCHOR.TOP,
    )

    _add_rect(slide, in_x, y + 3.10, box_w, 2.35, fill=COLORS["paper"], line=COLORS["border"], width=0.65, rounded=True)
    _add_text(slide, in_x + 0.12, y + 3.18, box_w - 0.24, 0.22, "2. Hierarchical structure map", size=8.5, bold=True, color=COLORS["green_dark"])
    _add_structure_table(slide, in_x + 0.12, y + 3.55, assets)

    _add_rect(slide, in_x, y + 6.05, box_w, 1.70, fill=COLORS["paper"], line=COLORS["border"], width=0.65, rounded=True)
    _add_text(slide, in_x + 0.12, y + 6.12, box_w - 0.24, 0.22, "3. Style reference", size=8.8, bold=True, color=COLORS["green_dark"])
    style_asset = _cover_asset(STYLE_REF, "figure1_style_reference.png", 520, 220)
    _add_picture_exact(slide, style_asset, in_x + 0.22, y + 6.48, box_w - 0.44, 1.05, background=RGBColor(255, 255, 255))

    model_x = x + 3.30
    _add_rect(slide, model_x, y + 3.10, 1.25, 3.35, fill=RGBColor(237, 245, 234), line=COLORS["green"], width=0.85, rounded=True)
    _add_text(slide, model_x + 0.12, y + 3.43, 1.00, 2.60, "SDXL\n+\nLoRA\n+\nHierarchical\nAdapter\n+\nStyle\nPrompting", size=11.0, bold=True, color=COLORS["text"], align=PP_ALIGN.CENTER)
    _add_line(slide, in_x + box_w, y + 1.82, model_x, y + 3.75, COLORS["green"], width=2.2, arrow=True)
    _add_line(slide, in_x + box_w, y + 4.18, model_x, y + 4.18, COLORS["green"], width=2.2, arrow=True)
    _add_line(slide, in_x + box_w, y + 6.92, model_x, y + 5.95, COLORS["green"], width=2.2, arrow=True)

    out_x = x + 4.96
    _add_text(slide, out_x - 0.05, y + 1.05, 2.35, 0.45, "Generated painting\n(Output)", size=10.2, bold=True, color=COLORS["green_dark"], align=PP_ALIGN.CENTER)
    output_asset = _cover_asset(OUTPUT_IMAGE, "figure1_output_square.png", 560, 560)
    _add_picture_exact(slide, output_asset, out_x - 0.08, y + 1.70, 2.28, 2.28, background=RGBColor(255, 255, 255))
    _add_line(slide, model_x + 1.25, y + 4.78, out_x - 0.08, y + 2.84, COLORS["green_dark"], width=2.8, arrow=True)
    checks = ["controllable composition", "meaningful blank space", "stable style", "high fidelity"]
    for idx, text in enumerate(checks):
        yy = y + 5.10 + idx * 0.55
        chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(out_x - 0.02), Inches(yy), Inches(0.30), Inches(0.30))
        chk.fill.solid()
        chk.fill.fore_color.rgb = COLORS["green_dark"]
        chk.line.color.rgb = COLORS["green_dark"]
        _add_text(slide, out_x + 0.02, yy + 0.02, 0.22, 0.22, "✓", size=10.2, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        _add_text(slide, out_x + 0.38, yy - 0.02, 2.00, 0.34, text, size=8.8, bold=True, color=COLORS["green_dark"])


def build() -> None:
    prs = _prepare_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    _add_panel_a(slide)
    _add_panel_b(slide)
    _add_panel_c(slide)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    build()
