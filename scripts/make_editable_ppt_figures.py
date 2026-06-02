from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_ROOT = ROOT / "data/processed/clp4k_v6_union_en"
ASSET_ROOT = ROOT / "paper/ppt_assets"
OUTPUT_PATH = ROOT / "paper/editable_figures_4_5.pptx"

FIGURE4_SAMPLES = [
    {
        "id": "clp4k_v6_union_en_000001",
        "prompt": "Misty mountains with river pavilion",
        "style": "ink wash",
        "labels": {
            "lora": "weak layout",
            "controlnet": "rigid contour",
            "ip_adapter": "style drift",
            "ours": "balanced output",
        },
    },
    {
        "id": "clp4k_v6_union_en_000110",
        "prompt": "Blue-green landscape with layered peaks",
        "style": "blue-green",
        "labels": {
            "lora": "loose hierarchy",
            "controlnet": "hard edges",
            "ip_adapter": "layout drift",
            "ours": "coherent style",
        },
    },
    {
        "id": "clp4k_v6_union_en_000061",
        "prompt": "Ink-wash river valley with blank space",
        "style": "freehand ink",
        "labels": {
            "lora": "weak blank",
            "controlnet": "stiff strokes",
            "ip_adapter": "space drift",
            "ours": "controlled blank",
        },
    },
    {
        "id": "clp4k_v6_union_en_000300",
        "prompt": "Light-red autumn distant mountains",
        "style": "light red",
        "labels": {
            "lora": "unstable tone",
            "controlnet": "rigid shape",
            "ip_adapter": "soft layout",
            "ours": "balanced tone",
        },
    },
]

FIGURE4_METHODS = [
    ("lora", "LoRA-only", ROOT / "outputs/paper_suite_v6/paper_lora_only_817/predictions_multiseed"),
    (
        "controlnet",
        "ControlNet",
        ROOT / "outputs/paper_suite_strong/paper_controlnet_817_seed42/predictions_multiseed",
    ),
    (
        "ip_adapter",
        "IP-Adapter only",
        ROOT / "outputs/paper_suite_strong/paper_ip_adapter_only_817_seed42/predictions_multiseed",
    ),
    ("ours", "Ours", ROOT / "outputs/paper_suite_v6/paper_ours_817/predictions_multiseed"),
]

UNIFIED_PROMPT = (
    "Chinese landscape painting; pine trees and waterside pavilion; "
    "misty distant mountains; ink wash; ample poetic blank space."
)


@dataclass(frozen=True)
class StructureCase:
    tag: str
    title: str
    sample_id: str
    blue_label: str
    orange_label: str
    blue_box: tuple[float, float, float, float]
    orange_box: tuple[float, float, float, float]


FIGURE5_CASES = [
    StructureCase(
        "A",
        "Large upper blank space",
        "clp4k_v6_union_en_000685",
        "blank space",
        "low composition",
        (0.06, 0.04, 0.82, 0.48),
        (0.18, 0.48, 0.92, 0.88),
    ),
    StructureCase(
        "B",
        "Dominant central peak",
        "clp4k_v6_union_en_000673",
        "spatial hierarchy",
        "main peak",
        (0.10, 0.08, 0.86, 0.88),
        (0.36, 0.08, 0.74, 0.82),
    ),
    StructureCase(
        "C",
        "Left-heavy mountain mass",
        "clp4k_v6_union_en_000740",
        "open right",
        "left mass",
        (0.42, 0.18, 0.92, 0.82),
        (0.06, 0.16, 0.48, 0.84),
    ),
    StructureCase(
        "D",
        "Lower-right foreground",
        "clp4k_v6_union_en_000001",
        "depth layers",
        "foreground",
        (0.08, 0.10, 0.74, 0.58),
        (0.58, 0.46, 0.92, 0.88),
    ),
]

FIGURE5_PREDICTION_ROOT = ROOT / "outputs/paper_suite_v6/paper_ours_817/predictions_multiseed"

WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(251, 250, 247)
TEXT = RGBColor(31, 41, 51)
MUTED = RGBColor(102, 112, 133)
BORDER = RGBColor(215, 210, 202)
BLUE = RGBColor(58, 141, 204)
ORANGE = RGBColor(198, 106, 46)
ACCENT = RGBColor(181, 51, 46)


def _structure_asset(sample_id: str) -> Path:
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    target = ASSET_ROOT / f"{sample_id}_lineart.png"
    if target.exists():
        return target
    source = DATA_ROOT / "structure_maps" / f"{sample_id}.png"
    rgba = Image.open(source).convert("RGBA")
    lineart = Image.fromarray(np.asarray(rgba)[..., 0]).convert("L")
    lineart = ImageOps.autocontrast(lineart)
    lineart.save(target)
    return target


def _add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 10,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(margin)
    shape.text_frame.margin_right = Inches(margin)
    shape.text_frame.margin_top = Inches(margin)
    shape.text_frame.margin_bottom = Inches(margin)
    shape.text_frame.vertical_anchor = valign
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor | None = WHITE,
    line: RGBColor = BORDER,
    width: float = 0.6,
    rounded: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def _set_dash(shape) -> None:
    line = shape._element.spPr.ln
    for existing in line.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash"):
        line.remove(existing)
    dash = line.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash", {"val": "dash"})
    line.append(dash)


def _add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, background: RGBColor = WHITE):
    _add_rect(slide, x, y, w, h, fill=background, line=BORDER, width=0.55)
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    picture_w = iw * scale
    picture_h = ih * scale
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    return slide.shapes.add_picture(str(path), Inches(picture_x), Inches(picture_y), Inches(picture_w), Inches(picture_h))


def _annotate_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    box: tuple[float, float, float, float],
    label: str,
    color: RGBColor,
) -> None:
    x0 = x + box[0] * w
    y0 = y + box[1] * h
    bw = (box[2] - box[0]) * w
    bh = (box[3] - box[1]) * h
    rect = _add_rect(slide, x0, y0, bw, bh, fill=None, line=color, width=1.1)
    _set_dash(rect)
    label_w = max(0.70, len(label) * 0.057)
    label_box = _add_rect(slide, x0 + 0.06, y0 + 0.05, label_w, 0.22, fill=WHITE, line=color, width=0.65, rounded=True)
    _add_textbox(slide, x0 + 0.09, y0 + 0.065, label_w - 0.06, 0.17, label, size=6.6, color=color)
    return label_box


def _prepare_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(12)
    return prs


def add_figure4_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    margin = 0.42
    gutter = 0.18
    header_y = 0.30
    table_y = 0.85
    row_h = 2.68
    condition_w = 4.30
    image_w = 3.53
    image_h = 2.18
    label_h = 0.28
    widths = [condition_w] + [image_w] * 4
    headers = ["Prompt / condition"] + [title for _, title, _ in FIGURE4_METHODS]

    x = margin
    for header, width in zip(headers, widths):
        _add_textbox(slide, x, header_y, width, 0.36, header, size=11.8, bold=True, align=PP_ALIGN.CENTER)
        x += width + gutter

    for row_index, sample in enumerate(FIGURE4_SAMPLES):
        y = table_y + row_index * row_h
        _add_rect(slide, margin, y, condition_w, row_h - 0.12, fill=WHITE, line=BORDER, width=0.65, rounded=True)
        _add_textbox(slide, margin + 0.18, y + 0.10, condition_w - 0.36, 0.48, sample["prompt"], size=9.4, bold=True)
        _add_textbox(slide, margin + 0.18, y + 0.52, condition_w - 0.36, 0.20, f"style: {sample['style']}", size=7.0, color=MUTED)
        structure = _structure_asset(sample["id"])
        _add_picture_contain(slide, structure, margin + 0.18, y + 0.82, condition_w - 0.36, 1.64, WHITE)

        x = margin + condition_w + gutter
        for method_key, _, directory in FIGURE4_METHODS:
            prediction = directory / f"{sample['id']}__seed42.png"
            _add_picture_contain(slide, prediction, x, y + 0.05, image_w, image_h, PAPER)
            _add_textbox(
                slide,
                x,
                y + image_h + 0.12,
                image_w,
                label_h,
                sample["labels"][method_key],
                size=8.3,
                bold=method_key == "ours",
                color=ACCENT if method_key == "ours" else MUTED,
                align=PP_ALIGN.CENTER,
            )
            x += image_w + gutter


def add_figure5_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    margin = 0.55
    _add_textbox(slide, margin, 0.18, 12.5, 0.40, "Figure 5. Structure Controllability Visualization", size=15, bold=True)
    _add_rect(slide, margin, 0.64, 18.90, 0.42, fill=RGBColor(247, 248, 250), line=BORDER, width=0.65, rounded=True)
    _add_textbox(slide, margin + 0.18, 0.70, 18.54, 0.24, f"Unified prompt: {UNIFIED_PROMPT}", size=9.0)

    left_x = margin
    left_w = 7.58
    right_x = 8.55
    right_w = 10.90
    _add_textbox(slide, left_x, 1.22, left_w, 0.30, "Input structure", size=10.8, bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, right_x, 1.22, right_w, 0.30, "Generated output (Ours)", size=10.8, bold=True, align=PP_ALIGN.CENTER)

    row_y = 1.60
    row_h = 2.50
    panel_w = 4.15
    panel_h = 1.82
    for index, case in enumerate(FIGURE5_CASES):
        y = row_y + index * row_h
        _add_rect(slide, left_x, y, left_w, row_h - 0.10, fill=PAPER, line=BORDER, width=0.65, rounded=True)
        _add_rect(slide, right_x, y, right_w, row_h - 0.10, fill=PAPER, line=BORDER, width=0.65, rounded=True)
        _add_textbox(slide, left_x + 0.20, y + 0.10, left_w - 0.40, 0.22, f"Case {case.tag}: {case.title}", size=8.7, bold=True)
        _add_textbox(slide, left_x + 0.20, y + 0.34, left_w - 0.40, 0.17, case.sample_id.replace("clp4k_v6_union_en_", "test #"), size=6.4, color=MUTED)

        structure_x = left_x + left_w - panel_w - 0.30
        structure_y = y + 0.54
        output_x = right_x + (right_w - panel_w) / 2
        output_y = y + 0.30
        structure = _structure_asset(case.sample_id)
        prediction = FIGURE5_PREDICTION_ROOT / f"{case.sample_id}__seed42.png"
        _add_picture_contain(slide, structure, structure_x, structure_y, panel_w, panel_h, WHITE)
        _add_picture_contain(slide, prediction, output_x, output_y, panel_w, panel_h, PAPER)
        for panel_x, panel_y in [(structure_x, structure_y), (output_x, output_y)]:
            _annotate_panel(slide, panel_x, panel_y, panel_w, panel_h, case.blue_box, case.blue_label, BLUE)
            _annotate_panel(slide, panel_x, panel_y, panel_w, panel_h, case.orange_box, case.orange_label, ORANGE)


def main() -> None:
    prs = _prepare_prs()
    add_figure4_slide(prs)
    add_figure5_slide(prs)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH.relative_to(ROOT)}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
